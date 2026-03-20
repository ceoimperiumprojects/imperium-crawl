#!/usr/bin/env python3
"""
imperium-crawl: Claude Code + MiniMax M2.7 LinkedIn Carousel Generator
Brand colors: #1A1A2E (bg), #FF2D78 (pink), #00D4FF (cyan), white text
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Brand Colors ──────────────────────────────────────
BG = RGBColor(0x1A, 0x1A, 0x2E)        # #1A1A2E
PINK = RGBColor(0xFF, 0x2D, 0x78)      # #FF2D78
CYAN = RGBColor(0x00, 0xD4, 0xFF)      # #00D4FF
WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # #FFFFFF
DARK_PINK = RGBColor(0xCC, 0x1F, 0x5E) # darker pink for depth
DARK_CYAN = RGBColor(0x00, 0xA8, 0xCC) # darker cyan

# ── Canvas: 10x10 inches (LinkedIn square optimal) ────
W = Inches(10)
H = Inches(10)

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=BG):
    """Fill entire slide with solid background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                font_name="Arial"):
    """Add a textbox with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_textbox_multiline(slide, left, top, width, height, lines, font_size=18,
                           default_color=WHITE, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, dict):
            text = line.get("text", "")
            color = line.get("color", default_color)
            size = line.get("size", font_size)
            bold = line.get("bold", False)
        else:
            text = line
            color = default_color
            size = font_size
            bold = False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])
    return txBox

def accent_bar(slide, color=PINK, width=Inches(0.15), height=Inches(0.6)):
    """Small vertical accent bar — used as bullet/decorator."""
    bar = add_rect(slide, Inches(0.5), Inches(2.2), width, height, color)
    return bar

# ───────────────────────────────────────────────────────
# SLIDE 1: Title
# ───────────────────────────────────────────────────────
def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Big gradient-like block at top (pink)
    add_rect(slide, 0, 0, W, Inches(3.5), PINK)

    # Decorative circles
    add_rect(slide, Inches(7.5), Inches(0.3), Inches(2), Inches(2), CYAN)
    add_rect(slide, Inches(8.5), Inches(1.8), Inches(1.2), Inches(1.2), DARK_PINK)

    # Title text in pink area
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(7), Inches(1.2),
                "Kako sam spojio", font_size=28, font_color=WHITE, bold=False)
    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(8), Inches(1.4),
                "Claude Code + MiniMax M2.7", font_size=38, font_color=WHITE, bold=True)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(7), Inches(0.8),
                "za autonomous AI improvement loop", font_size=22, font_color=BG)

    # Main tagline below pink block
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.8),
                "Od Haiku besplatnog → Haiku za $0.03/M tokens", font_size=24,
                font_color=CYAN, bold=True)

    # imperium-crawl badge
    add_rect(slide, Inches(0.5), Inches(5.2), Inches(3.5), Inches(0.7), DARK_PINK)
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(3.5), Inches(0.7),
                "  imperium-crawl  ", font_size=16, font_color=WHITE, bold=True,
                align=PP_ALIGN.CENTER)

    # Subtext
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(9), Inches(2.5),
                "Full stack setup • LiteLLM proxy • Overnight autoresearch",
                font_size=16, font_color=RGBColor(0xAA, 0xAA, 0xBB))

    # Bottom accent
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), CYAN)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 2: Problem (CCR broken)
# ───────────────────────────────────────────────────────
def slide_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Header bar
    add_rect(slide, 0, 0, W, Inches(1.2), DARK_PINK)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                "PROBLEM", font_size=16, font_color=PINK, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.7),
                "Claude Code Router mi je uništio billing", font_size=26, font_color=WHITE, bold=True)

    # Problem cards
    problems = [
        (PINK, "💳 Max subscription → API Usage Billing",
         "CCR hardkodira auth token = Anthropic ne vidi više Max sub. Svaki request košta."),
        (CYAN, "💀 Router padne → Claude Code mrtav",
         "Nema fallbacka. Nema recovery-ja. Samo timeout i crni ekran."),
        (RGBColor(0xFF, 0x88, 0x00), "🔒 OAuth token se gubi u tranzitu",
         "CCR nikad ne prosleđuje realni OAuth header iz Claude Code klijenta."),
    ]

    for i, (color, title, desc) in enumerate(problems):
        top = Inches(1.5) + i * Inches(2.3)
        add_rect(slide, Inches(0.5), top, Inches(0.15), Inches(1.8), color)
        add_rect(slide, Inches(0.75), top, Inches(8.7), Inches(1.8),
                 RGBColor(0x22, 0x22, 0x3E))
        add_textbox(slide, Inches(1.0), top + Inches(0.15), Inches(8.2), Inches(0.6),
                    title, font_size=18, font_color=color, bold=True)
        add_textbox(slide, Inches(1.0), top + Inches(0.75), Inches(8.2), Inches(0.9),
                    desc, font_size=14, font_color=WHITE)

    # Bottom bar
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), PINK)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 3: Solution (LiteLLM)
# ───────────────────────────────────────────────────────
def slide_solution(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Header
    add_rect(slide, 0, 0, W, Inches(1.2), DARK_CYAN)
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                "REŠENJE", font_size=16, font_color=CYAN, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.7),
                "LiteLLM proxy — transparentna forward proxy", font_size=26, font_color=WHITE, bold=True)

    # Key insight box
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5), RGBColor(0x22, 0x22, 0x3E))
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(0.15), Inches(1.5), CYAN)
    add_textbox(slide, Inches(0.85), Inches(1.6), Inches(8.4), Inches(0.6),
                "Jedna settings linija rešava sve:", font_size=16, font_color=CYAN)
    add_textbox(slide, Inches(0.85), Inches(2.15), Inches(8.4), Inches(0.6),
                "forward_client_headers_to_llm_api: true", font_size=20, font_color=WHITE, bold=True)

    # How it works
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(9), Inches(0.6),
                "Kako radi:", font_size=18, font_color=PINK, bold=True)

    steps = [
        ("1", "Claude Code šalje OAuth header sa svakim requestom"),
        ("2", "LiteLLM proxy prima request, prosleđuje headers dalje"),
        ("3", "Anthropic vidi originalni OAuth → Max subscription ✅"),
    ]

    for i, (num, text) in enumerate(steps):
        top = Inches(3.9) + i * Inches(1.1)
        add_rect(slide, Inches(0.5), top, Inches(0.6), Inches(0.6), CYAN)
        add_textbox(slide, Inches(0.5), top, Inches(0.6), Inches(0.6),
                    num, font_size=20, font_color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.3), top + Inches(0.05), Inches(8.2), Inches(0.6),
                    text, font_size=16, font_color=WHITE)

    # Bottom accent
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), CYAN)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 4: Architecture diagram
# ───────────────────────────────────────────────────────
def slide_architecture(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Header
    add_rect(slide, 0, 0, W, Inches(1.2), RGBColor(0x15, 0x15, 0x28))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                "ARHITEKTURA", font_size=16, font_color=PINK, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.7),
                "localhost:4000 — tačno na sred", font_size=26, font_color=WHITE, bold=True)

    # ── Layer 1: Claude Code ──
    cc_y = Inches(2.0)
    add_rect(slide, Inches(0.5), cc_y, Inches(2.5), Inches(1.2),
             RGBColor(0x33, 0x33, 0x55), line_color=PINK, line_width=Pt(2))
    add_textbox(slide, Inches(0.5), cc_y + Inches(0.1), Inches(2.5), Inches(0.5),
                "Claude Code", font_size=14, font_color=PINK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), cc_y + Inches(0.55), Inches(2.5), Inches(0.5),
                "haiku / sonnet / opus", font_size=11, font_color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow down to LiteLLM
    add_rect(slide, Inches(1.65), cc_y + Inches(1.2), Inches(0.1), Inches(0.7), PINK)

    # ── Layer 2: LiteLLM ──
    ll_y = Inches(3.6)
    add_rect(slide, Inches(2.5), ll_y, Inches(5), Inches(1.4),
             RGBColor(0x00, 0x44, 0x66), line_color=CYAN, line_width=Pt(3))
    add_textbox(slide, Inches(2.5), ll_y + Inches(0.1), Inches(5), Inches(0.5),
                "LiteLLM Proxy", font_size=18, font_color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), ll_y + Inches(0.6), Inches(5), Inches(0.5),
                "localhost:4000", font_size=13, font_color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), ll_y + Inches(1.0), Inches(5), Inches(0.4),
                "forward_client_headers: true", font_size=10,
                font_color=RGBColor(0x88, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

    # Arrows down from LiteLLM
    # Left arrow (sonnet/opus → Anthropic)
    add_rect(slide, Inches(2.0), ll_y + Inches(1.4), Inches(0.1), Inches(0.8), CYAN)
    add_rect(slide, Inches(0.5), ll_y + Inches(2.1), Inches(2.2), Inches(0.1), CYAN)
    add_rect(slide, Inches(0.5), ll_y + Inches(2.1), Inches(0.1), Inches(0.8), CYAN)

    # Right arrow (haiku → MiniMax)
    add_rect(slide, Inches(7.4), ll_y + Inches(1.4), Inches(0.1), Inches(0.8), PINK)
    add_rect(slide, Inches(7.4), ll_y + Inches(2.1), Inches(2.2), Inches(0.1), PINK)
    add_rect(slide, Inches(9.5), ll_y + Inches(2.1), Inches(0.1), Inches(0.8), PINK)

    # ── Layer 3A: Anthropic ──
    anth_y = ll_y + Inches(3.0)
    add_rect(slide, Inches(0.5), anth_y, Inches(2.2), Inches(1.2),
             RGBColor(0x44, 0x22, 0x55), line_color=PINK, line_width=Pt(2))
    add_textbox(slide, Inches(0.5), anth_y + Inches(0.1), Inches(2.2), Inches(0.5),
                "Anthropic", font_size=14, font_color=PINK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), anth_y + Inches(0.55), Inches(2.2), Inches(0.5),
                "sonnet / opus", font_size=11, font_color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), anth_y + Inches(0.9), Inches(2.2), Inches(0.3),
                "OAuth ✅ Max sub", font_size=9, font_color=CYAN, align=PP_ALIGN.CENTER)

    # ── Layer 3B: MiniMax ──
    add_rect(slide, Inches(7.4), anth_y, Inches(2.2), Inches(1.2),
             RGBColor(0x33, 0x00, 0x22), line_color=PINK, line_width=Pt(2))
    add_textbox(slide, Inches(7.4), anth_y + Inches(0.1), Inches(2.2), Inches(0.5),
                "MiniMax M2.7", font_size=14, font_color=PINK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.4), anth_y + Inches(0.55), Inches(2.2), Inches(0.5),
                "haiku", font_size=11, font_color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.4), anth_y + Inches(0.9), Inches(2.2), Inches(0.3),
                "$0.03/M tokens 💰", font_size=9, font_color=CYAN, align=PP_ALIGN.CENTER)

    # Labels on arrows
    add_textbox(slide, Inches(0.1), ll_y + Inches(1.7), Inches(1.5), Inches(0.5),
                "OAuth", font_size=12, font_color=CYAN)
    add_textbox(slide, Inches(7.8), ll_y + Inches(1.7), Inches(1.5), Inches(0.5),
                "MiniMax API", font_size=12, font_color=PINK)

    # Bottom
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), PINK)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 5: Results
# ───────────────────────────────────────────────────────
def slide_results(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Header
    add_rect(slide, 0, 0, W, Inches(1.2), RGBColor(0x15, 0x15, 0x28))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                "REZULTAT", font_size=16, font_color=CYAN, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.7),
                "Sve radi, ništa ne košta, ništa ne puca", font_size=24, font_color=WHITE, bold=True)

    # Stats row
    stats = [
        (CYAN, "50x", "jeftinije", "Haiku na MiniMax vs Max API"),
        (PINK, "$0", "sonnet/opus", "Max subscription — FREE"),
        (RGBColor(0x00, 0xCC, 0x66), "∞", "rezilijentnost", "Fallback ako LiteLLM padne"),
    ]

    for i, (color, big, label, sub) in enumerate(stats):
        left = Inches(0.5) + i * Inches(3.1)
        add_rect(slide, left, Inches(1.5), Inches(2.9), Inches(2.5),
                 RGBColor(0x22, 0x22, 0x3E))
        add_rect(slide, left, Inches(1.5), Inches(2.9), Inches(0.1), color)
        add_textbox(slide, left, Inches(1.7), Inches(2.9), Inches(1.0),
                    big, font_size=44, font_color=color, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left, Inches(2.75), Inches(2.9), Inches(0.5),
                    label, font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left, Inches(3.25), Inches(2.9), Inches(0.6),
                    sub, font_size=11, font_color=RGBColor(0xAA, 0xAA, 0xBB), align=PP_ALIGN.CENTER)

    # Resilience point
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.8), RGBColor(0x22, 0x22, 0x3E))
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(0.15), Inches(0.8), RGBColor(0x00, 0xCC, 0x66))
    add_textbox(slide, Inches(0.85), Inches(4.4), Inches(8.4), Inches(0.6),
                "✅  Healthcheck u bashrc — ako LiteLLM padne, Claude Code koristi direktni Anthropic.",
                font_size=14, font_color=WHITE)

    # Bottom accent
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), PINK)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 6: Autoresearch Loop
# ───────────────────────────────────────────────────────
def slide_loop(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Header
    add_rect(slide, 0, 0, W, Inches(1.2), RGBColor(0x15, 0x15, 0x28))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                "AUTORESEARCH LOOP", font_size=16, font_color=PINK, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.7),
                "Karpathy-style autonomous improvement", font_size=24, font_color=WHITE, bold=True)

    # Loop steps — circular/timeline layout
    steps = [
        (CYAN,    "1", "Eval", "Run test suite → find lowest score"),
        (PINK,    "2", "Fix", "Make ONE targeted improvement"),
        (RGBColor(0x00, 0xCC, 0x66), "3", "Test", "Re-run evaluation"),
        (RGBColor(0xFF, 0x88, 0x00), "4", "Commit", "If improved → git commit. If not → discard."),
        (CYAN,    "5", "Repeat", "Next iteration — overnight, 50x"),
    ]

    for i, (color, num, title, desc) in enumerate(steps):
        top = Inches(1.5) + i * Inches(1.35)
        # Circle number
        add_rect(slide, Inches(0.5), top, Inches(0.7), Inches(0.7), color)
        add_textbox(slide, Inches(0.5), top, Inches(0.7), Inches(0.7),
                    num, font_size=22, font_color=BG, bold=True, align=PP_ALIGN.CENTER)
        # Connector line (except last)
        if i < len(steps) - 1:
            add_rect(slide, Inches(0.83), top + Inches(0.7), Inches(0.04), Inches(0.6), color)
        # Content
        add_textbox(slide, Inches(1.4), top + Inches(0.05), Inches(2.5), Inches(0.5),
                    title, font_size=20, font_color=color, bold=True)
        add_textbox(slide, Inches(1.4), top + Inches(0.5), Inches(8), Inches(0.6),
                    desc, font_size=14, font_color=WHITE)

    # Bottom bar
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), PINK)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 7: Impact
# ───────────────────────────────────────────────────────
def slide_impact(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Header
    add_rect(slide, 0, 0, W, Inches(1.2), RGBColor(0x15, 0x15, 0x28))
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.7),
                "IMPACT", font_size=16, font_color=CYAN, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.55), Inches(9), Inches(0.7),
                "Šta ovo znači za OSS?", font_size=28, font_color=WHITE, bold=True)

    # Main quote/insight
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.8), RGBColor(0x22, 0x22, 0x3E))
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(0.15), Inches(1.8), PINK)
    add_textbox(slide, Inches(0.85), Inches(1.7), Inches(8.4), Inches(0.7),
                '"The future of OSS maintenance might just be', font_size=18, font_color=WHITE, bold=True)
    add_textbox(slide, Inches(0.85), Inches(2.4), Inches(8.4), Inches(0.7),
                'leave it running overnight."', font_size=22, font_color=CYAN, bold=True)

    # Key points
    points = [
        ("🚀", "imperium-crawl", "AI-native web scraping framework, 30+ tools, stalno se sam poboljšava"),
        ("💡", "Compound growth", "Male izmene se akumuliraju — svake noći malo bolji scoring"),
        ("🔧", "Minimal cost", "50 iteracija noću košta ~$0.50 na MiniMax umesto ~$25 na Max API"),
        ("⚡", "Max sub preserved", "sonnet/opus i dalje besplatni sa Max — nikad se ne troše na haiku"),
    ]

    for i, (emoji, title, desc) in enumerate(points):
        top = Inches(3.5) + i * Inches(1.35)
        add_rect(slide, Inches(0.5), top, Inches(9), Inches(1.15),
                 RGBColor(0x22, 0x22, 0x3E))
        add_textbox(slide, Inches(0.7), top + Inches(0.1), Inches(1.2), Inches(0.5),
                    emoji, font_size=22)
        add_textbox(slide, Inches(1.9), top + Inches(0.1), Inches(7.4), Inches(0.45),
                    title, font_size=16, font_color=PINK, bold=True)
        add_textbox(slide, Inches(1.9), top + Inches(0.55), Inches(7.4), Inches(0.5),
                    desc, font_size=13, font_color=WHITE)

    # Bottom
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), CYAN)

    return slide

# ───────────────────────────────────────────────────────
# SLIDE 8: CTA
# ───────────────────────────────────────────────────────
def slide_cta(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Big CTA area
    add_rect(slide, 0, 0, W, Inches(6), PINK)

    # Decorative elements
    add_rect(slide, Inches(7.5), Inches(4), Inches(2.5), Inches(2), DARK_PINK)
    add_rect(slide, Inches(8.5), Inches(0.5), Inches(1.2), Inches(1.2), CYAN)
    add_rect(slide, Inches(0.5), Inches(5), Inches(1.5), Inches(1), CYAN)

    # Main CTA
    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(0.8),
                "imperium-crawl", font_size=22, font_color=WHITE)
    add_textbox(slide, Inches(0.5), Inches(1.9), Inches(8.5), Inches(1.2),
                "Open source. 30+ scraping tools.", font_size=32, font_color=WHITE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(8.5), Inches(0.8),
                "Autonomous AI improvement loop included.", font_size=20, font_color=BG)
    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(8.5), Inches(0.8),
                "Try it tonight.", font_size=24, font_color=BG, bold=True)

    # GitHub reference
    add_rect(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(1.0), RGBColor(0x22, 0x22, 0x3E))
    add_rect(slide, Inches(0.5), Inches(6.2), Inches(0.15), Inches(1.0), CYAN)
    add_textbox(slide, Inches(0.85), Inches(6.3), Inches(8.4), Inches(0.8),
                "github.com/[username]/imperium-crawl  ← link u opisu",
                font_size=14, font_color=CYAN)

    # Question
    add_textbox(slide, Inches(0.5), Inches(7.5), Inches(9), Inches(1.5),
                "Running AI agents at scale?", font_size=20, font_color=WHITE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(8.2), Inches(9), Inches(1.0),
                "Kako rešavaš model routing i cost optimization? 🔥",
                font_size=18, font_color=CYAN)

    # Bottom bar
    add_rect(slide, 0, Inches(9.7), W, Inches(0.3), CYAN)

    return slide

# ───────────────────────────────────────────────────────
# MAIN — Generate and save
# ───────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_results(prs)
    slide_loop(prs)
    slide_impact(prs)
    slide_cta(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "carousel-claude-code-minimax.pptx")
    prs.save(out_path)

    size_kb = os.path.getsize(out_path) // 1024
    print(f"✅ Saved: {out_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"💾 File size: {size_kb} KB")

if __name__ == "__main__":
    main()
